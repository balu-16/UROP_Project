from pathlib import Path

from fastapi import HTTPException, UploadFile, status


class DocumentParser:
    supported_extensions = {".pdf", ".txt", ".md", ".markdown", ".pptx"}

    async def parse_upload(self, file: UploadFile, max_mb: int) -> tuple[str, dict]:
        filename = file.filename or "document.txt"
        suffix = Path(filename).suffix.lower()
        if suffix not in self.supported_extensions:
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"Unsupported file type: {suffix}")
        data = await file.read()
        if len(data) > max_mb * 1024 * 1024:
            raise HTTPException(status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE, detail="File is too large")
        if suffix == ".pdf":
            text = self._parse_pdf(data)
        elif suffix == ".pptx":
            text = self._parse_pptx(data)
        else:
            text = data.decode("utf-8", errors="ignore")
        return text, {"source": filename, "extension": suffix, "size_bytes": len(data)}

    def _parse_pdf(self, data: bytes) -> str:
        try:
            from pypdf import PdfReader
            from io import BytesIO

            reader = PdfReader(BytesIO(data))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception as exc:
            raise HTTPException(status_code=400, detail=f"Unable to parse PDF: {exc}") from exc

    def _parse_pptx(self, data: bytes) -> str:
        try:
            from io import BytesIO

            from pptx import Presentation

            prs = Presentation(BytesIO(data))
            slides_text: list[str] = []
            for idx, slide in enumerate(prs.slides, start=1):
                parts: list[str] = []
                for shape in slide.shapes:
                    try:
                        if shape.has_text_frame:
                            t = (shape.text or "").strip()
                            if t:
                                parts.append(t)
                        if shape.has_table:
                            for row in shape.table.rows:
                                cells = [(cell.text or "").strip() for cell in row.cells]
                                cells = [c for c in cells if c]
                                if cells:
                                    parts.append(" | ".join(cells))
                    except Exception:
                        continue
                try:
                    notes = (slide.notes_slide.placeholders[1].text or "").strip() if slide.has_notes_slide else ""
                    if notes:
                        parts.append(f"Notes: {notes}")
                except Exception:
                    pass
                slide_text = "\n".join(parts).strip()
                if slide_text:
                    slides_text.append(f"[Slide {idx}]\n{slide_text}")
            return "\n\n".join(slides_text)
        except Exception as exc:
            raise HTTPException(status_code=400, detail=f"Unable to parse PPTX: {exc}") from exc

